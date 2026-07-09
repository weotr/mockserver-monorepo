#!/usr/bin/env python3
"""Build + render the ARCHITECTURE-style AI diagrams (AI overview, contract
verification, centralised deployment) in the MockServer flat house style — white
boxes with black borders and blue arrows, like
`system_under_test_with_mockserver_proxy.png`.

Unlike the flow diagrams (render_diagrams.py reads their slide geometry), these
have no flow-slide source. So each diagram is defined ONCE as a box/arrow spec
below, and this script emits BOTH:
  * a matching editable slide appended to MockServerScenarios.pptx (Blank layout),
    so the layout can be tweaked in PowerPoint, and
  * the PNG the website ships,
from the same coordinates — slide and image stay in sync.

Idempotent: drops any slide it previously appended (index >= APPEND_BASELINE)
before re-adding. Run AFTER build_ai_diagrams.py (which appends the flow slides).

Usage:  python3 render_ai_architecture.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont

HERE = Path(__file__).resolve().parent
PPTX = HERE.parent / "MockServerScenarios.pptx"
IMAGES = HERE.parent
APPEND_BASELINE = 23  # 19 committed + 4 flow slides from build_ai_diagrams.py

BLUE = (0x4F, 0x81, 0xBD)
BLACK = (0x40, 0x40, 0x40)
GREY = (0x80, 0x80, 0x80)
FONTS = ["/System/Library/Fonts/Supplemental/Arial.ttf",
         "/System/Library/Fonts/Helvetica.ttc"]

# ---- diagram specs: canvas is 0..CW x 0..CH; boxes/arrows in those units ----
# box: (id, x, y, w, h, "label", emphasis)   arrow: (ax, ay, bx, by, "label")
# caption: (x, y, "text")  — small grey section label, left-aligned

def col(s, x, w, label, emph=False): return dict(id=s, x=x, y=0, w=w, h=0, label=label, emph=emph)

OVERVIEW = dict(
    name="mockserver_ai_overview", CW=1500, CH=760,
    boxes=[
        ("a1", 60, 70, 300, 110, "AI Assistant\n(Claude Code, Cursor)", False),
        ("a2", 760, 70, 320, 110, "MockServer", True),
        ("b1", 60, 320, 300, 110, "Application\nunder test", False),
        ("b2", 760, 320, 360, 110, "MockServer\n(fake LLM / MCP / A2A)", True),
        ("c1", 60, 580, 260, 110, "AI Agent", False),
        ("c2", 620, 580, 320, 110, "MockServer\n(transparent proxy)", True),
        ("c3", 1180, 580, 260, 110, "Real LLM API", False),
    ],
    arrows=[
        ("a1", "a2", "drives via MCP tools"),
        ("b1", "b2", "HTTP / SSE / WebSocket"),
        ("c1", "c2", "via proxy"),
        ("c2", "c3", "forwards"),
    ],
    captions=[(60, 30, "Mode 1 - AI controls MockServer"),
              (60, 280, "Mode 2 - MockServer mocks AI services"),
              (60, 540, "Mode 3 - MockServer observes AI traffic")],
)

VERIFICATION = dict(
    name="mockserver_ai_verification", CW=1500, CH=760,
    boxes=[
        ("p1", 60, 70, 320, 110, "Recorded traffic\n(MockServer log)", False),
        ("p2", 620, 70, 360, 110, "Verify vs OpenAPI", True),
        ("p3", 1180, 70, 280, 110, "Conformance\nreport", False),
        ("a1", 60, 320, 280, 110, "OpenAPI spec", False),
        ("a2", 560, 320, 360, 110, "Run contract test", True),
        ("a3", 1140, 320, 300, 110, "Live service", False),
        ("r1", 60, 580, 280, 110, "OpenAPI spec", False),
        ("r2", 560, 580, 360, 110, "Run resiliency test", True),
        ("r3", 1140, 580, 300, 110, "Live service", False),
    ],
    arrows=[
        ("p1", "p2", "recorded pairs"), ("p2", "p3", "per-pair result"),
        ("a1", "a2", ""), ("a2", "a3", "example requests"),
        ("r1", "r2", ""), ("r2", "r3", "malformed inputs"),
    ],
    captions=[(60, 30, "1 - Passive validation (recorded traffic)"),
              (60, 280, "2 - Active contract test (live service)"),
              (60, 540, "3 - Resiliency / fuzz test (live service)")],
)

CENTRAL_SINGLE = dict(
    name="mockserver_centralized_single", CW=1500, CH=620,
    boxes=[
        ("t1", 60, 60, 260, 100, "Team A CI", False),
        ("t2", 60, 250, 260, 100, "Team B CI", False),
        ("t3", 60, 440, 260, 100, "Team C CI", False),
        ("lb", 560, 250, 300, 100, "Load Balancer\n(optional)", False),
        ("ms", 1080, 250, 360, 100, "MockServer\n(single shared node)", True),
    ],
    arrows=[("t1", "lb", ""), ("t2", "lb", ""), ("t3", "lb", ""), ("lb", "ms", "")],
    captions=[(60, 20, "Single shared instance")],
)

CENTRAL_CLUSTER = dict(
    name="mockserver_centralized_clustered", CW=1560, CH=680,
    boxes=[
        ("t1", 40, 120, 240, 100, "Team A CI", False),
        ("t2", 40, 360, 240, 100, "Team B CI", False),
        ("lb", 470, 240, 280, 100, "Load Balancer", False),
        ("na", 980, 40, 300, 100, "MockServer Node A", True),
        ("nb", 980, 280, 300, 100, "MockServer Node B", True),
        ("nc", 980, 520, 300, 100, "MockServer Node C", True),
        ("inf", 980, 640, 300, 0, "", False),  # placeholder for caption row; see captions
    ],
    arrows=[("t1", "lb", ""), ("t2", "lb", ""),
            ("lb", "na", ""), ("lb", "nb", ""), ("lb", "nc", "")],
    captions=[(40, 70, "Clustered HA - shared state via Infinispan (JGroups REPL_SYNC)")],
    replication=["na", "nb", "nc"],  # draw a bidirectional spine linking these nodes
)

DIAGRAMS = [OVERVIEW, VERIFICATION, CENTRAL_SINGLE, CENTRAL_CLUSTER]

# ----------------------------- PNG rendering -----------------------------
def font(px):
    for f in FONTS:
        try:
            return ImageFont.truetype(f, px)
        except Exception:
            pass
    return ImageFont.load_default()

def boxmap(spec):
    return {b[0]: b for b in spec["boxes"] if b[3] and b[4]}

def edge_point(b, toward):
    x, y, w, h = b[1], b[2], b[3], b[4]
    cx, cy = x + w / 2, y + h / 2
    tx, ty = toward
    if abs(tx - cx) >= abs(ty - cy):
        return (x + w if tx > cx else x, cy)
    return (cx, y + h if ty > cy else y)

def render_png(spec):
    SS = 2
    W, H = spec["CW"] * SS, spec["CH"] * SS
    img = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(img)
    fb, fl, fc = font(22 * SS), font(19 * SS), font(20 * SS)
    bm = boxmap(spec)
    # arrows first (under boxes)
    for ax, bx, label in spec["arrows"]:
        a, b = bm[ax], bm[bx]
        ca = (a[1] + a[3] / 2, a[2] + a[4] / 2)
        cb = (b[1] + b[3] / 2, b[2] + b[4] / 2)
        p = edge_point(a, cb); q = edge_point(b, ca)
        P = (p[0] * SS, p[1] * SS); Q = (q[0] * SS, q[1] * SS)
        d.line([P, Q], fill=BLUE, width=4 * SS)
        # arrowhead at Q
        import math
        ang = math.atan2(Q[1] - P[1], Q[0] - P[0]); s = 16 * SS
        d.polygon([Q, (Q[0] - s * math.cos(ang - 0.5), Q[1] - s * math.sin(ang - 0.5)),
                   (Q[0] - s * math.cos(ang + 0.5), Q[1] - s * math.sin(ang + 0.5))], fill=BLUE)
        if label:
            mx, my = (P[0] + Q[0]) / 2, (P[1] + Q[1]) / 2
            tb = d.textbbox((0, 0), label, font=fl)
            d.text((mx - (tb[2] - tb[0]) / 2, my - (tb[3] - tb[1]) - 8 * SS), label, fill=BLACK, font=fl)
    # replication spine (clustered)
    if spec.get("replication"):
        ids = spec["replication"]; xs = [bm[i] for i in ids]
        x = (xs[0][1]) * SS - 28 * SS
        top = (xs[0][2] + xs[0][4] / 2) * SS; bot = (xs[-1][2] + xs[-1][4] / 2) * SS
        d.line([(x, top), (x, bot)], fill=BLUE, width=4 * SS)
        for b in xs:
            yy = (b[2] + b[4] / 2) * SS
            d.line([(x, yy), (b[1] * SS, yy)], fill=BLUE, width=4 * SS)
    # boxes
    for bid, x, y, w, h, label, emph in spec["boxes"]:
        if not (w and h):
            continue
        bb = (x * SS, y * SS, (x + w) * SS, (y + h) * SS)
        d.rectangle(bb, fill="white", outline=(BLUE if emph else BLACK), width=(5 if emph else 3) * SS)
        lines = label.split("\n")
        th = sum(d.textbbox((0, 0), ln, font=fb)[3] for ln in lines)
        cy = (bb[1] + bb[3]) / 2 - th / 2
        for ln in lines:
            tb = d.textbbox((0, 0), ln, font=fb)
            d.text(((bb[0] + bb[2]) / 2 - (tb[2] - tb[0]) / 2, cy), ln, fill=BLACK, font=fb)
            cy += tb[3]
    # captions
    for cx, cy, text in spec["captions"]:
        d.text((cx * SS, cy * SS), text, fill=GREY, font=fc)
    img = img.resize((spec["CW"], spec["CH"]), Image.LANCZOS)
    dst = IMAGES / (spec["name"] + ".png")
    img.save(dst)
    print("rendered", dst.name, img.size)

# ----------------------------- pptx slide -----------------------------
def add_arrowhead(connector):
    ln = connector.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)

def emu_x(spec, v, sw, m):  return Emu(int(m + v / spec["CW"] * (sw - 2 * m)))
def emu_y(spec, v, sh, m):  return Emu(int(m + v / spec["CH"] * (sh - 2 * m)))

def build_slide(prs, spec):
    blank = next(l for l in prs.slide_layouts if l.name == "Blank")
    s = prs.slides.add_slide(blank)
    sw, sh, m = prs.slide_width, prs.slide_height, Emu(380000)
    sx = lambda v: emu_x(spec, v, sw, m); sy = lambda v: emu_y(spec, v, sh, m)
    bm = boxmap(spec)
    for ax, bx, label in spec["arrows"]:
        a, b = bm[ax], bm[bx]
        ca = (a[1] + a[3] / 2, a[2] + a[4] / 2); cb = (b[1] + b[3] / 2, b[2] + b[4] / 2)
        p = edge_point(a, cb); q = edge_point(b, ca)
        cxn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx(p[0]), sy(p[1]), sx(q[0]), sy(q[1]))
        cxn.line.color.rgb = RGBColor(*BLUE); cxn.line.width = Pt(2.25)
        add_arrowhead(cxn)
    for bid, x, y, w, h, label, emph in spec["boxes"]:
        if not (w and h):
            continue
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx(x), sy(y),
                                Emu(int(w / spec["CW"] * (sw - 2 * m))), Emu(int(h / spec["CH"] * (sh - 2 * m))))
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        sp.line.color.rgb = RGBColor(*(BLUE if emph else BLACK)); sp.line.width = Pt(2.5 if emph else 1.5)
        sp.shadow.inherit = False
        tf = sp.text_frame; tf.word_wrap = True
        tf.text = label.replace("\n", "\v")
        for para in tf.paragraphs:
            para.alignment = 2  # center
            for r in para.runs:
                r.font.size = Pt(13); r.font.color.rgb = RGBColor(*BLACK)
    for cx, cy, text in spec["captions"]:
        tb = s.shapes.add_textbox(sx(cx), sy(cy), Emu(6000000), Emu(380000))
        tb.text_frame.text = text
        r = tb.text_frame.paragraphs[0].runs[0]; r.font.size = Pt(12); r.font.color.rgb = RGBColor(*GREY)
    return s

def main():
    prs = Presentation(str(PPTX))
    lst = prs.slides._sldIdLst
    for sldId in list(lst)[APPEND_BASELINE:]:
        prs.part.drop_rel(sldId.get(qn('r:id')))
        lst.remove(sldId)
    for spec in DIAGRAMS:
        print("Architecture diagram:", spec["name"])
        render_png(spec)
        build_slide(prs, spec)
    prs.save(str(PPTX))
    n = len(list(prs.slides._sldIdLst))
    print(f"\nSaved {PPTX.name}: {n} slides; appended {len(DIAGRAMS)} architecture slides at idx {APPEND_BASELINE}..{n-1}.")

if __name__ == '__main__':
    main()
