#!/usr/bin/env python3
"""Append the AI / load-injection FLOW diagrams to MockServerScenarios.pptx as
NEW editable slides, by cloning existing house-style flow slides and relabelling
them. This is the ADDITIVE companion to build_diagrams.py.

Why a separate, additive script: the committed deck has 19 slides — idx 0..14 are
genuine sources (incl. the hand-made TLS, single-page-app and retrieve-logs
slides) and idx 15..18 are previously generated flow PNgrams. build_diagrams.py
still assumes ORIGINAL_SLIDES = 11 and would DELETE slides 11..18 if run, so we do
NOT touch it. This script only ever appends, and is idempotent: on each run it
drops any slide at index >= BASELINE and re-adds the definitions below.

Flow diagrams render to PNG from slide geometry via render_diagrams.py. The
ARCHITECTURE-style diagrams (AI overview, verification, centralised deployment)
are not flow slides and are produced by render_ai_architecture.py, which also adds
their matching editable slides — keep those out of here.

Usage:  python3 build_ai_diagrams.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn

# reuse the cloning / relabel / recolour helpers from the original builder
from build_diagrams import duplicate_slide, relabel, recolor

PPTX = Path(__file__).resolve().parent.parent / "MockServerScenarios.pptx"
BASELINE = 19  # committed slide count; everything appended by this script lives here onward

# Each entry: (source slide idx to clone, relabel map, optional accent to recolour)
FLOW_DIAGRAMS = [
    # Load injection — drive load from a scenario, collect SLO samples (clone the
    # proxy+matcher+cylinder optimise layout, idx 10).
    (10, {
        'Proxy': ['Load', 'Scenario'],
        'Request/Matcher': ['Virtual', 'Users'],
        'LLM/Traffic': ['SLO', 'Samples'],
        '1. Proxy LLM Call': ['1. Trigger Run'],
        '2. Record Traffic': ['2. Send to Target'],
        '3. Request Brief': ['3. Ramp Virtual Users'],
        '4. Analyse Traffic': ['4. Record Metrics'],
        '5. Export Brief': ['5. Verify SLO'],
    }, 'accent3'),  # green-ish, distinct from proxy/optimise teal
    # AI traffic inspection — proxy + record + view (clone record/proxy, idx 9).
    (9, {
        'Expectation/or/Proxy': ['Proxy'],
        'Recorded/Requests': ['Captured', 'Traffic'],
        'Request/Matcher': ['Dashboard', '& Retrieve'],
        '1. Receive Request': ['1. Proxy AI Call'],
        '2. Record Request': ['2. Record Traffic'],
        '3. Send Verification': ['3. Open Dashboard'],
        '4. Match Request(s)': ['4. Inspect Traffic'],
        '5. Return Result': ['5. View / Export'],
    }, None),
    # AI debugging — forward, record, retrieve, diagnose (clone record/proxy, idx 9).
    (9, {
        'Expectation/or/Proxy': ['Proxy'],
        'Recorded/Requests': ['Recorded', 'Traffic'],
        'Request/Matcher': ['AI', 'Assistant'],
        '1. Receive Request': ['1. Proxy Request'],
        '2. Record Request': ['2. Record Exchange'],
        '3. Send Verification': ['3. Retrieve Pair'],
        '4. Match Request(s)': ['4. Diagnose'],
        '5. Return Result': ['5. Explain Fix'],
    }, 'accent6'),  # orange, distinct from traffic inspection (avoid red = faults)
    # AI protocol mocking — receive, match, return scripted response (clone the LLM
    # response action, idx 4).
    (4, {
        'LLM/Response': ['MCP / A2A', 'Response'],
        '1. Receive Prompt': ['1. Receive MCP Call'],
        '2. Match Request': ['2. Match Request'],
        '3. Stream Tokens': ['3. Stream SSE Result'],
    }, 'accent4'),  # purple, like the LLM-response mock diagram (returns a mock response)
]

def main():
    prs = Presentation(str(PPTX))
    lst = prs.slides._sldIdLst
    # idempotent: drop anything previously appended by this script (index >= BASELINE)
    for sldId in list(lst)[BASELINE:]:
        prs.part.drop_rel(sldId.get(qn('r:id')))
        lst.remove(sldId)

    for i, (src, mapping, accent) in enumerate(FLOW_DIAGRAMS):
        print(f"Flow diagram {i} (clone slide idx {src})")
        s = duplicate_slide(prs, src)
        relabel(s, mapping)
        if accent:
            recolor(s, accent)

    prs.save(str(PPTX))
    n = len(list(prs.slides._sldIdLst))
    print(f"\nSaved {PPTX.name}: {n} slides ({BASELINE} kept + {n - BASELINE} appended flow diagrams).")
    print("New flow slide indices:", list(range(BASELINE, n)))

if __name__ == '__main__':
    main()
